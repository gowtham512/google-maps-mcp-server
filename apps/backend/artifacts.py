"""Generate downloadable artifacts (PPTX/PDF/JSON) from structured artifact data."""

import io
import json
import re
from typing import Any

import httpx
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import Image as RLImage
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle


ARTIFACT_FORMATS = {"json", "pptx", "pdf"}

# Max image download size to avoid memory issues (10 MB).
MAX_IMAGE_BYTES = 10 * 1024 * 1024
# Timeout for fetching remote images.
IMAGE_TIMEOUT = 10.0


def sanitize_filename(title: str) -> str:
    """Turn an artifact title into a safe file-name slug."""
    slug = re.sub(r"[^\w\s-]", "", title).strip().lower()
    slug = re.sub(r"[-\s]+", "-", slug)
    return slug or "artifact"


def _parse_artifact(raw: str | None) -> dict[str, Any] | None:
    if not raw:
        return None
    try:
        data = json.loads(raw)
    except json.JSONDecodeError:
        return None
    if not isinstance(data, dict):
        return None
    return data


def load_artifact(data: str | None) -> dict[str, Any] | None:
    """Validate and normalize stored artifact JSON."""
    artifact = _parse_artifact(data)
    if not artifact:
        return None
    artifact_type = artifact.get("type")
    if artifact_type not in {"slides", "report"}:
        return None
    artifact.setdefault("title", "Artifact")
    if artifact_type == "slides":
        artifact.setdefault("slides", [])
    elif artifact_type == "report":
        artifact.setdefault("sections", [])
    return artifact


def resolve_format(artifact_type: str | None, format: str) -> str:
    """Resolve 'auto' to the natural format for an artifact type."""
    if format != "auto":
        return format
    return "pptx" if artifact_type == "slides" else "pdf"


async def _fetch_image(url: str) -> bytes | None:
    """Download an image from a URL, returning None on any failure."""
    if not url or not url.startswith(("http://", "https://")):
        return None
    try:
        async with httpx.AsyncClient(timeout=IMAGE_TIMEOUT, follow_redirects=True) as client:
            resp = await client.get(url)
            if resp.status_code != 200:
                return None
            content = resp.content
            if len(content) > MAX_IMAGE_BYTES:
                return None
            # Basic content-type validation.
            content_type = resp.headers.get("content-type", "")
            if not content_type.startswith(("image/", "application/octet-stream")):
                return None
            return content
    except Exception:
        return None


def build_json(artifact: dict[str, Any]) -> bytes:
    """Return the artifact as pretty-printed JSON bytes."""
    return json.dumps(artifact, indent=2, ensure_ascii=False).encode("utf-8")


def _add_image_to_slide(slide, image_bytes: bytes, left: float, top: float, width: float, height: float) -> None:
    """Embed image bytes into a pptx slide at the given position/size."""
    try:
        image_stream = io.BytesIO(image_bytes)
        slide.shapes.add_picture(image_stream, Inches(left), Inches(top), Inches(width), Inches(height))
    except Exception:
        # If the image is invalid, ignore it.
        pass


async def build_pptx(artifact: dict[str, Any]) -> bytes:
    """Build an editable PowerPoint from a slides artifact."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Title slide
    slide = prs.slides.add_slide(blank_layout)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = artifact.get("title", "Slides")
    p.font.size = Pt(44)
    p.font.bold = True

    slides = artifact.get("slides", [])
    for s in slides:
        slide = prs.slides.add_slide(blank_layout)

        image_url = s.get("image_url")
        image_bytes = await _fetch_image(image_url) if image_url else None

        # Layout: if there's an image, put it on the right half and text on the left.
        if image_bytes:
            text_left = 0.5
            text_width = 5.8
            _add_image_to_slide(slide, image_bytes, 6.8, 1.5, 5.5, 4.5)
        else:
            text_left = 0.5
            text_width = 12.333

        header = slide.shapes.add_textbox(Inches(text_left), Inches(0.4), Inches(text_width), Inches(1.2))
        tf = header.text_frame
        p = tf.paragraphs[0]
        p.text = s.get("title", "")
        p.font.size = Pt(36)
        p.font.bold = True

        if s.get("subtitle"):
            sub_top = 1.6
            sub = slide.shapes.add_textbox(Inches(text_left), Inches(sub_top), Inches(text_width), Inches(0.6))
            tf = sub.text_frame
            p = tf.paragraphs[0]
            p.text = s["subtitle"]
            p.font.size = Pt(22)

        body_top = Inches(2.4) if s.get("subtitle") else Inches(1.8)
        body = slide.shapes.add_textbox(Inches(text_left), body_top, Inches(text_width), Inches(4.6))
        tf = body.text_frame
        tf.word_wrap = True

        bullets = s.get("bullets", [])
        if bullets:
            for idx, bullet in enumerate(bullets):
                if idx == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"• {bullet}"
                p.font.size = Pt(22)
                p.space_after = Pt(12)
        else:
            p = tf.paragraphs[0]
            p.text = "No content."
            p.font.size = Pt(20)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


async def build_pdf(artifact: dict[str, Any]) -> bytes:
    """Build a printable PDF from a report artifact using ReportLab."""
    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=letter, rightMargin=54, leftMargin=54, topMargin=54, bottomMargin=54
    )
    styles = getSampleStyleSheet()
    story: list[Any] = []

    title_style = styles["Heading1"]
    story.append(Paragraph(artifact.get("title", "Report"), title_style))
    story.append(Spacer(1, 0.2 * inch))

    for section in artifact.get("sections", []):
        heading = section.get("heading", "")
        body = section.get("body", "")
        image_url = section.get("image_url")
        image_bytes = await _fetch_image(image_url) if image_url else None

        if heading:
            story.append(Paragraph(heading, styles["Heading2"]))

        if image_bytes:
            try:
                image_stream = io.BytesIO(image_bytes)
                img = RLImage(image_stream, width=5 * inch, height=3 * inch)
                img.hAlign = "LEFT"
                story.append(img)
                story.append(Spacer(1, 0.1 * inch))
            except Exception:
                # If the image can't be rendered, skip it.
                pass

        if body:
            story.append(Paragraph(body.replace("\n", "<br/>"), styles["Normal"]))
        story.append(Spacer(1, 0.2 * inch))

    if not artifact.get("sections"):
        story.append(Paragraph("No content.", styles["Normal"]))

    doc.build(story)
    buf.seek(0)
    return buf.getvalue()
